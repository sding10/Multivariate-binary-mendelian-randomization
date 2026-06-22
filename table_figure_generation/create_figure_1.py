# ============================================================
# Create Figure 1a, 1b, 1c from raw CSV summary file
# using true min / Q1 / median / Q3 / max values
# ============================================================

import re
from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle, Patch
from matplotlib.ticker import FixedLocator, FixedFormatter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# File paths
# ============================================================

# The user said the CSV is inside the Data folder in the current working folder.
CSV_PATH = Path("Data") / "Figure1_Bias_Boxplot_Summary.csv"

OUTPUT_DIR = Path("Figure1_bias_from_csv_output")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

PPTX_OUT = OUTPUT_DIR / "Figure1_bias_from_csv.pptx"


# ============================================================
# Figure configuration
# ============================================================

FIGURE_CONFIG = {
    "1a": {"conf": 0.5},
    "1b": {"conf": 1.0},
    "1c": {"conf": 1.5},
}


# ============================================================
# Plot settings & styles
# ============================================================

Y_LIMITS_REAL = (-10, 10)
Y_TICKS_REAL = [-10, -5, -2, -1, 0, 1, 2, 5, 10]

METHOD_ORDER = ["2SPS", "2SRI", "GMM", "IV-MVB"]

# Grayscale + muted dark blue, consistent with your recent version
METHOD_FILL = {
    "2SPS": "#F0F0F0",
    "2SRI": "#C2C2C2",
    "GMM": "#7A7A7A",
    "IV-MVB": "#1C3144",
}

METHOD_EDGE = {
    "2SPS": "#444444",
    "2SRI": "#444444",
    "GMM": "#333333",
    "IV-MVB": "#1C3144",
}

IV_STRENGTH_ORDER = [
    "Very Weak IVs",
    "Weak IVs",
    "Moderate IVs",
    "Strong IVs",
    "Very Strong IVs",
]

SAMPLE_SIZE_ORDER = [500, 1500, 2500, 5000, 7500, 10000]

plt.rcParams.update({
    "font.family": "Arial",
    "font.size": 11,
    "axes.labelsize": 12,
    "axes.titlesize": 12,
    "xtick.labelsize": 10.5,
    "ytick.labelsize": 10.5,
    "figure.dpi": 150,
    "savefig.dpi": 600,
    "axes.linewidth": 1.0,
})


# ============================================================
# Data loading and filtering
# ============================================================

def clean_x_label(value):
    """
    Wrap IV-strength labels neatly for presentation.
    """
    text = str(value)

    replacements = {
        "Very Weak IVs": "Very Weak\nIVs",
        "Weak IVs": "Weak\nIVs",
        "Moderate IVs": "Moderate\nIVs",
        "Strong IVs": "Strong\nIVs",
        "Very Strong IVs": "Very Strong\nIVs",
    }

    return replacements.get(text, text)


def load_and_filter_raw_csv(csv_path):
    """
    Load the raw summary CSV and keep only the rows needed for Figure 1.

    Required filtering:
      - ignore rows where b1 is 2 or 3
      - keep:
          A with b1 = 0
          B with b1 = 1
          C with b1 = 0
          D with b1 = 1
    """
    df = pd.read_csv(csv_path)

    expected_cols = {
        "scenario", "b1", "n", "c_x", "iv_strength", "method",
        "min_bias", "q1_bias", "median_bias", "q3_bias", "max_bias"
    }
    missing = expected_cols - set(df.columns)
    if missing:
        raise ValueError(f"CSV is missing required columns: {missing}")

    # Keep only the scenario-specific b1 values we want
    mask = (
        ((df["scenario"] == "A") & (df["b1"] == 0)) |
        ((df["scenario"] == "B") & (df["b1"] == 1)) |
        ((df["scenario"] == "C") & (df["b1"] == 0)) |
        ((df["scenario"] == "D") & (df["b1"] == 1))
    )

    df = df.loc[mask].copy()

    # Keep only the methods in desired order
    df = df[df["method"].isin(METHOD_ORDER)].copy()

    # Enforce method order
    df["method"] = pd.Categorical(df["method"], categories=METHOD_ORDER, ordered=True)

    return df


def get_figure_data(df_all, conf_strength):
    """
    Subset one confounding strength (c_x == c_y).
    """
    df = df_all[np.isclose(df_all["c_x"], conf_strength)].copy()

    # Helpful plotting x-value:
    # Scenarios A/B vary iv_strength
    # Scenarios C/D vary n
    def choose_x(row):
        if row["scenario"] in ["A", "B"]:
            return row["iv_strength"]
        else:
            return int(row["n"])

    df["x_value"] = df.apply(choose_x, axis=1)

    return df


# ============================================================
# Rendering helpers
# ============================================================

def clip_to_axis(y):
    """
    Clip any y-value to the visible plotting limits.
    """
    return max(min(y, Y_LIMITS_REAL[1]), Y_LIMITS_REAL[0])


def visible_interval(y0, y1):
    """
    Return the visible part of [y0, y1] after clipping to axis limits.
    Return None if no visible interval remains.
    """
    low = max(min(y0, y1), Y_LIMITS_REAL[0])
    high = min(max(y0, y1), Y_LIMITS_REAL[1])

    if high <= low:
        return None
    return low, high


def draw_box_from_summary(ax, x_center, min_v, q1, med, q3, max_v, method, box_width):
    """
    Draw one boxplot-like object directly from precomputed summary values:
      min, Q1, median, Q3, max
    All visible geometry is clipped to [-10, 10].
    """
    fill = METHOD_FILL[method]
    edge = METHOD_EDGE[method]

    # ------------------------------
    # Box: Q1 to Q3
    # ------------------------------
    box_interval = visible_interval(q1, q3)
    if box_interval is not None:
        box_low, box_high = box_interval
        rect = Rectangle(
            (x_center - box_width / 2, box_low),
            box_width,
            box_high - box_low,
            facecolor=fill,
            edgecolor=edge,
            linewidth=1.0,
            zorder=3,
        )
        ax.add_patch(rect)

    # ------------------------------
    # Median line
    # ------------------------------
    if Y_LIMITS_REAL[0] <= med <= Y_LIMITS_REAL[1]:
        ax.plot(
            [x_center - box_width / 2, x_center + box_width / 2],
            [med, med],
            color="#E41A1C",
            linewidth=2.4,
            zorder=5,
        )

    # ------------------------------
    # Lower whisker: min to Q1
    # ------------------------------
    lower_whisker = visible_interval(min_v, q1)
    if lower_whisker is not None:
        low, high = lower_whisker
        ax.plot([x_center, x_center], [low, high], color=edge, linewidth=1.0, zorder=2)
        ax.plot(
            [x_center - box_width * 0.30, x_center + box_width * 0.30],
            [low, low],
            color=edge,
            linewidth=1.0,
            zorder=2,
        )

    # ------------------------------
    # Upper whisker: Q3 to max
    # ------------------------------
    upper_whisker = visible_interval(q3, max_v)
    if upper_whisker is not None:
        low, high = upper_whisker
        ax.plot([x_center, x_center], [low, high], color=edge, linewidth=1.0, zorder=2)
        ax.plot(
            [x_center - box_width * 0.30, x_center + box_width * 0.30],
            [high, high],
            color=edge,
            linewidth=1.0,
            zorder=2,
        )


def plot_panel(ax, df, scenario, panel_label):
    """
    Plot one scenario panel.
    """
    panel = df[df["scenario"] == scenario].copy()

    if scenario in ["A", "B"]:
        x_order = IV_STRENGTH_ORDER
    else:
        x_order = SAMPLE_SIZE_ORDER

    x_positions = np.arange(len(x_order))
    group_width = 0.65
    box_width = 0.13
    offsets = np.linspace(-group_width / 2, group_width / 2, len(METHOD_ORDER))

    for i, x_val in enumerate(x_order):
        for method, offset in zip(METHOD_ORDER, offsets):
            row = panel[(panel["x_value"] == x_val) & (panel["method"] == method)]
            if row.empty:
                continue

            row = row.iloc[0]

            draw_box_from_summary(
                ax=ax,
                x_center=i + offset,
                min_v=row["min_bias"],
                q1=row["q1_bias"],
                med=row["median_bias"],
                q3=row["q3_bias"],
                max_v=row["max_bias"],
                method=method,
                box_width=box_width,
            )

    # Symlog axis retained per your request
    ax.set_yscale("symlog", linthresh=1.0, linscale=1.0, base=10)

    # Reference lines
    ax.axhline(0, color="#2CA02C", linewidth=1.0, zorder=1)
    ax.axhline(1, color="#E2E2E2", linewidth=0.8, linestyle="--", zorder=1)
    ax.axhline(-1, color="#E2E2E2", linewidth=0.8, linestyle="--", zorder=1)

    ax.set_xlim(-0.7, len(x_order) - 0.3)
    ax.set_ylim(Y_LIMITS_REAL)

    # Force your preferred y-ticks
    ax.yaxis.set_major_locator(FixedLocator(Y_TICKS_REAL))
    ax.yaxis.set_major_formatter(FixedFormatter([str(t) for t in Y_TICKS_REAL]))

    ax.set_xticks(x_positions)
    ax.set_xticklabels([clean_x_label(x) for x in x_order])

    ax.set_ylabel("Estimated Bias in Causal Parameter", fontweight="bold")

    ax.set_facecolor("white")
    for spine in ax.spines.values():
        spine.set_linewidth(1.0)
        spine.set_color("#222222")

    ax.tick_params(axis="both", direction="out", length=5, width=1.0, colors="black")

    # Panel label
    ax.text(
        0.02, 0.97, panel_label,
        transform=ax.transAxes,
        ha="left",
        va="top",
        fontsize=20,
        fontweight="bold",
        color="black",
    )


def create_one_figure(df, fig_label, conf_strength):
    """
    Create one complete 2x2 figure.
    Slightly more rectangular than a square.
    """
    fig, axes = plt.subplots(
        2, 2,
        figsize=(11.2, 8.7),
        facecolor="white"
    )
    axes = axes.ravel()

    plot_panel(axes[0], df, "A", "A")
    plot_panel(axes[1], df, "B", "B")
    plot_panel(axes[2], df, "C", "C")
    plot_panel(axes[3], df, "D", "D")

    # Layout tuned for clean whitespace and a consolidated legend
    fig.subplots_adjust(
        left=0.10,
        right=0.96,
        top=0.96,
        bottom=0.10,
        wspace=0.24,
        hspace=0.30,
    )

    # Consolidated legend at bottom
    legend_handles = [
        Patch(facecolor=METHOD_FILL[m], edgecolor=METHOD_EDGE[m], label=m, linewidth=1.0)
        for m in METHOD_ORDER
    ]
    leg = fig.legend(
        handles=legend_handles,
        loc="lower center",
        ncol=4,
        frameon=False,
        handlelength=1.2,
        handletextpad=0.4,
        columnspacing=2.0,
        fontsize=11,
    )
    for text in leg.get_texts():
        text.set_fontweight("bold")

    conf_tag = str(conf_strength).replace(".", "p")

    png_path = OUTPUT_DIR / f"Figure_{fig_label}_bias_c{conf_tag}.png"
    svg_path = OUTPUT_DIR / f"Figure_{fig_label}_bias_c{conf_tag}.svg"

    fig.savefig(
        png_path,
        dpi=600,
        facecolor="white",
        edgecolor="white",
        bbox_inches="tight",
    )

    fig.savefig(
        svg_path,
        facecolor="white",
        edgecolor="white",
        bbox_inches="tight",
    )

    plt.close(fig)

    return png_path, svg_path


# ============================================================
# PowerPoint helpers
# ============================================================

def make_figure_title(fig_label, conf):
    """
    Slide title and subtitle.
    """
    title = f"Figure {fig_label}. Bias distributions across four simulation scenarios with 30 instruments"
    subtitle = (
        f"Confounding strength: c_x = c_y = {conf:.1f}. "
        f"Scenarios A/B vary instrument strength; Scenarios C/D vary sample size. "
        f"See Table 1 for more details on the parameter settings in each scenario."
    )
    return title, subtitle


def add_slide_decorations(slide, title, subtitle):
    """
    Add title and subtitle to the slide.
    """
    title_box = slide.shapes.add_textbox(
        Inches(0.45), Inches(0.18), Inches(14.0), Inches(0.70)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 20, 20)

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.45), Inches(0.65), Inches(14.0), Inches(0.70)
    )
    tf2 = subtitle_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.name = "Arial"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(60, 60, 60)


def create_powerpoint(image_records):
    """
    Create the PowerPoint with one slide per figure.
    """
    prs = Presentation()
    prs.slide_width = Inches(15.0)
    prs.slide_height = Inches(11.5)

    blank_layout = prs.slide_layouts[6]

    for rec in image_records:
        slide = prs.slides.add_slide(blank_layout)

        title, subtitle = make_figure_title(rec["fig_label"], rec["conf"])
        add_slide_decorations(slide, title, subtitle)

        # Slightly rectangular image footprint, left aligned,
        # with whitespace left on the right side of the slide.
        img_left = Inches(1.35)
        img_top = Inches(1.78)
        img_width = Inches(10.75)
        img_height = Inches(8.35)

        # A lighter border than before
        border_pad = Inches(0.04)
        border_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            img_left - border_pad,
            img_top - border_pad,
            img_width + 2 * border_pad,
            img_height + 2 * border_pad,
        )
        border_box.fill.background()
        border_box.line.color.rgb = RGBColor(55, 55, 55)
        border_box.line.width = Pt(1.5)

        slide.shapes.add_picture(
            str(rec["png"]),
            img_left,
            img_top,
            width=img_width,
            height=img_height,
        )

    prs.save(PPTX_OUT)


# ============================================================
# Main
# ============================================================

def main():
    df_all = load_and_filter_raw_csv(CSV_PATH)

    image_records = []

    for fig_label, cfg in FIGURE_CONFIG.items():
        conf = cfg["conf"]
        df_fig = get_figure_data(df_all, conf)

        png_path, svg_path = create_one_figure(
            df=df_fig,
            fig_label=fig_label,
            conf_strength=conf,
        )

        image_records.append({
            "fig_label": fig_label,
            "conf": conf,
            "png": png_path,
            "svg": svg_path,
        })

    create_powerpoint(image_records)

    print("\nSuccess! Figure 1 PowerPoint created successfully.\n")
    print(f"Output PPTX: {PPTX_OUT}\n")
    for rec in image_records:
        print(f"Figure {rec['fig_label']} PNG: {rec['png']}")
        print(f"Figure {rec['fig_label']} SVG: {rec['svg']}")


if __name__ == "__main__":
    main()