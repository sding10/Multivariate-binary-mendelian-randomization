# ============================================================
# Create Figure 2 from summary CSV file
#
# Figure 2. Bias distributions by number of instrumental variables
# under Scenario B (weak IVs)
#
# Scenario B:
#   beta1 = 1
#   N = 1,000
#   c_x = c_y = 1.5
#   weak IVs
#
# The CSV contains summary estimates:
#   n_iv, method, median, Q1, Q3, min, max, numconverged
# ============================================================

from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle, Patch
from matplotlib.ticker import FixedLocator, FixedFormatter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# File paths
# ============================================================

CSV_NAME = "Paper3_NumIV_Weak_bias_summary_by_niv_method.csv"

CSV_PATH = Path("Data") / CSV_NAME

# Fallback: useful if the CSV is in the current working folder
if not CSV_PATH.exists():
    CSV_PATH = Path(CSV_NAME)

OUTPUT_DIR = Path("Figure2_bias_by_num_iv_output")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

PPTX_OUT = OUTPUT_DIR / "Figure2_bias_by_num_iv.pptx"


# ============================================================
# Plot settings, matched to Figure 1
# ============================================================

Y_LIMITS_REAL = (-10, 10)
Y_TICKS_REAL = [-10, -5, -2, -1, 0, 1, 2, 5, 10]

METHOD_ORDER = ["2SPS", "2SRI", "GMM", "IV-MVB"]

METHOD_FILL = {
    "2SPS": "#F0F0F0",
    "2SRI": "#C2C2C2",
    "GMM": "#7A7A7A",
    "IV-MVB": "#1C3144",
}

METHOD_EDGE = {
    "2SPS": "#444444",
    "2SRI": "#444444",
    "GMM": "#333333",
    "IV-MVB": "#1C3144",
}

plt.rcParams.update({
    "font.family": "Arial",
    "font.size": 11,
    "axes.labelsize": 12,
    "axes.titlesize": 12,
    "xtick.labelsize": 10.5,
    "ytick.labelsize": 10.5,
    "figure.dpi": 150,
    "savefig.dpi": 600,
    "axes.linewidth": 1.0,
})


# ============================================================
# Data loading
# ============================================================

def load_summary_csv(csv_path):
    df = pd.read_csv(csv_path)

    expected_cols = {
        "n_iv",
        "method",
        "median",
        "Q1",
        "Q3",
        "min",
        "max",
        "numconverged",
    }

    missing = expected_cols - set(df.columns)
    if missing:
        raise ValueError(f"CSV is missing required columns: {missing}")

    df = df[df["method"].isin(METHOD_ORDER)].copy()

    df["method"] = pd.Categorical(
        df["method"],
        categories=METHOD_ORDER,
        ordered=True,
    )

    df["n_iv"] = df["n_iv"].astype(int)

    return df


# ============================================================
# Rendering helpers
# ============================================================

def visible_interval(y0, y1):
    """
    Return the visible part of [y0, y1] after clipping to axis limits.
    Return None if no visible interval remains.
    """
    low = max(min(y0, y1), Y_LIMITS_REAL[0])
    high = min(max(y0, y1), Y_LIMITS_REAL[1])

    if high <= low:
        return None

    return low, high


def draw_box_from_summary(
    ax,
    x_center,
    min_v,
    q1,
    med,
    q3,
    max_v,
    method,
    box_width,
):
    """
    Draw one boxplot-like object directly from summary values:
      min, Q1, median, Q3, max

    The visible geometry is clipped to [-10, 10], matching Figure 1.
    """
    fill = METHOD_FILL[method]
    edge = METHOD_EDGE[method]

    # Box: Q1 to Q3
    box_interval = visible_interval(q1, q3)
    if box_interval is not None:
        box_low, box_high = box_interval

        rect = Rectangle(
            (x_center - box_width / 2, box_low),
            box_width,
            box_high - box_low,
            facecolor=fill,
            edgecolor=edge,
            linewidth=1.0,
            zorder=3,
        )
        ax.add_patch(rect)

    # Median line
    if Y_LIMITS_REAL[0] <= med <= Y_LIMITS_REAL[1]:
        ax.plot(
            [x_center - box_width / 2, x_center + box_width / 2],
            [med, med],
            color="#E41A1C",
            linewidth=2.4,
            zorder=5,
        )

    # Lower whisker: min to Q1
    lower_whisker = visible_interval(min_v, q1)
    if lower_whisker is not None:
        low, high = lower_whisker

        ax.plot(
            [x_center, x_center],
            [low, high],
            color=edge,
            linewidth=1.0,
            zorder=2,
        )

        ax.plot(
            [x_center - box_width * 0.30, x_center + box_width * 0.30],
            [low, low],
            color=edge,
            linewidth=1.0,
            zorder=2,
        )

    # Upper whisker: Q3 to max
    upper_whisker = visible_interval(q3, max_v)
    if upper_whisker is not None:
        low, high = upper_whisker

        ax.plot(
            [x_center, x_center],
            [low, high],
            color=edge,
            linewidth=1.0,
            zorder=2,
        )

        ax.plot(
            [x_center - box_width * 0.30, x_center + box_width * 0.30],
            [high, high],
            color=edge,
            linewidth=1.0,
            zorder=2,
        )


# ============================================================
# Plotting
# ============================================================

def create_figure2(df):
    n_iv_order = sorted(df["n_iv"].unique())

    fig, ax = plt.subplots(
        1,
        1,
        figsize=(11.2, 6.6),
        facecolor="white",
    )

    x_positions = np.arange(len(n_iv_order))

    group_width = 0.65
    box_width = 0.13
    offsets = np.linspace(-group_width / 2, group_width / 2, len(METHOD_ORDER))

    for i, n_iv in enumerate(n_iv_order):
        for method, offset in zip(METHOD_ORDER, offsets):
            row = df[
                (df["n_iv"] == n_iv) &
                (df["method"] == method)
            ]

            if row.empty:
                continue

            row = row.iloc[0]

            draw_box_from_summary(
                ax=ax,
                x_center=i + offset,
                min_v=row["min"],
                q1=row["Q1"],
                med=row["median"],
                q3=row["Q3"],
                max_v=row["max"],
                method=method,
                box_width=box_width,
            )

    # Symlog axis, matched to Figure 1
    ax.set_yscale("symlog", linthresh=1.0, linscale=1.0, base=10)

    # Reference lines, matched to Figure 1
    ax.axhline(0, color="#2CA02C", linewidth=1.0, zorder=1)
    ax.axhline(1, color="#E2E2E2", linewidth=0.8, linestyle="--", zorder=1)
    ax.axhline(-1, color="#E2E2E2", linewidth=0.8, linestyle="--", zorder=1)

    ax.set_xlim(-0.7, len(n_iv_order) - 0.3)
    ax.set_ylim(Y_LIMITS_REAL)

    ax.yaxis.set_major_locator(FixedLocator(Y_TICKS_REAL))
    ax.yaxis.set_major_formatter(FixedFormatter([str(t) for t in Y_TICKS_REAL]))

    ax.set_xticks(x_positions)
    ax.set_xticklabels([str(x) for x in n_iv_order])

    ax.set_xlabel("Number of Instrumental Variables", fontweight="bold")
    ax.set_ylabel("Estimated Bias in Causal Parameter", fontweight="bold")

    ax.set_facecolor("white")

    for spine in ax.spines.values():
        spine.set_linewidth(1.0)
        spine.set_color("#222222")

    ax.tick_params(
        axis="both",
        direction="out",
        length=5,
        width=1.0,
        colors="black",
    )

    ax.grid(
        axis="y",
        color="#E6E6E6",
        linewidth=0.8,
        linestyle="-",
        zorder=0,
    )

    legend_handles = [
        Patch(
            facecolor=METHOD_FILL[m],
            edgecolor=METHOD_EDGE[m],
            label=m,
            linewidth=1.0,
        )
        for m in METHOD_ORDER
    ]

    leg = fig.legend(
        handles=legend_handles,
        loc="lower center",
        ncol=4,
        frameon=False,
        handlelength=1.2,
        handletextpad=0.4,
        columnspacing=2.0,
        fontsize=11,
    )

    for text in leg.get_texts():
        text.set_fontweight("bold")

    fig.subplots_adjust(
        left=0.10,
        right=0.96,
        top=0.94,
        bottom=0.18,
    )

    png_path = OUTPUT_DIR / "Figure2_bias_by_num_iv.png"
    svg_path = OUTPUT_DIR / "Figure2_bias_by_num_iv.svg"

    fig.savefig(
        png_path,
        dpi=600,
        facecolor="white",
        edgecolor="white",
        bbox_inches="tight",
    )

    fig.savefig(
        svg_path,
        facecolor="white",
        edgecolor="white",
        bbox_inches="tight",
    )

    plt.close(fig)

    return png_path, svg_path


# ============================================================
# PowerPoint helpers
# ============================================================

def make_figure_title():
    title = (
        "Figure 2. Bias distributions by number of instrumental variables\n"
        "under Scenario B (weak IVs)"
    )

    subtitle = (
        "Confounding strength: cₓ = cᵧ = 1.5. "
        "Scenario B (β₁ = 1, N = 1,000) with weak IVs. "
        "Each number-of-IVs setting was simulated independently. "
        "See Table 1 for additional design parameters."
    )

    return title, subtitle


def add_slide_decorations(slide, title, subtitle):
    title_box = slide.shapes.add_textbox(
        Inches(0.45),
        Inches(0.18),
        Inches(14.0),
        Inches(0.82),
    )

    tf = title_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 20, 20)

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.45),
        Inches(0.88),
        Inches(14.0),
        Inches(0.75),
    )

    tf2 = subtitle_box.text_frame
    tf2.word_wrap = True

    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.name = "Arial"
    p2.font.size = Pt(13.5)
    p2.font.color.rgb = RGBColor(60, 60, 60)


def create_powerpoint(png_path):
    prs = Presentation()
    prs.slide_width = Inches(15.0)
    prs.slide_height = Inches(11.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    title, subtitle = make_figure_title()
    add_slide_decorations(slide, title, subtitle)

    # Consistent with Figure 1 slide positioning:
    # image centered but slightly left-of-center, leaving more whitespace on right.
    img_left = Inches(1.35)
    img_top = Inches(1.92)
    img_width = Inches(10.75)
    img_height = Inches(7.15)

    border_pad = Inches(0.04)

    border_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        img_left - border_pad,
        img_top - border_pad,
        img_width + 2 * border_pad,
        img_height + 2 * border_pad,
    )

    border_box.fill.background()
    border_box.line.color.rgb = RGBColor(55, 55, 55)
    border_box.line.width = Pt(1.5)

    slide.shapes.add_picture(
        str(png_path),
        img_left,
        img_top,
        width=img_width,
        height=img_height,
    )

    prs.save(PPTX_OUT)


# ============================================================
# Main
# ============================================================

def main():
    df = load_summary_csv(CSV_PATH)

    png_path, svg_path = create_figure2(df)
    create_powerpoint(png_path)

    print("\nSuccess! Figure 2 PowerPoint created successfully.\n")
    print(f"Input CSV: {CSV_PATH}")
    print(f"Output PPTX: {PPTX_OUT}")
    print(f"Figure 2 PNG: {png_path}")
    print(f"Figure 2 SVG: {svg_path}")


if __name__ == "__main__":
    main()