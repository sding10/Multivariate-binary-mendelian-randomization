# ============================================================
# Create Figure 5
# Comparison of MR method performance across confounding strengths
# under weak instruments
# ============================================================

import re
from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle
from matplotlib.ticker import FixedLocator, FixedFormatter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# File paths
# ============================================================

BIAS_CSV_PATH = Path("Data") / "Figure1_Bias_Boxplot_Summary.csv"
COVERAGE_CSV_PATH = Path("Data") / "Figure_Coverage_Summary.csv"
TYPE1_CSV_PATH = Path("Data") / "Figure_TypeIError_Summary.csv"
POWER_CSV_PATH = Path("Data") / "Figure_Power_Summary.csv"

OUTPUT_DIR = Path("Figure5_operating_characteristics_output")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

PPTX_OUT = OUTPUT_DIR / "Figure5_operating_characteristics.pptx"


# ============================================================
# Fixed Figure 5 settings
# ============================================================

TARGET_N = 1000
TARGET_IV_STRENGTH = "weak"
CONFOUNDING_ORDER = [0.5, 1.0, 1.5]
CONFOUNDING_LABELS = ["0.5", "1.0", "1.5"]

SCENARIO_NULL = "A"
SCENARIO_POWER = "B"

B1_NULL = 0
B1_POWER = 1


# ============================================================
# Plot settings
# ============================================================

METHOD_ORDER = ["2SPS", "2SRI", "GMM", "IV-MVB"]

METHOD_FILL = {
    "2SPS": "#F0F0F0",
    "2SRI": "#C2C2C2",
    "GMM": "#7A7A7A",
    "IV-MVB": "#1C3144",
}

METHOD_EDGE = {
    "2SPS": "#444444",
    "2SRI": "#444444",
    "GMM": "#333333",
    "IV-MVB": "#1C3144",
}

METHOD_COLOR = {
    "2SPS": "#CFCFCF",
    "2SRI": "#AFAFAF",
    "GMM": "#8E8E8E",
    "IV-MVB": "#111111",
}

METHOD_MARKER = {
    "2SPS": "o",
    "2SRI": "*",
    "GMM": "D",
    "IV-MVB": "s",
}

METHOD_LINESTYLE = {
    "2SPS": (0, (10, 4)),
    "2SRI": (0, (10, 4)),
    "GMM": (0, (10, 4)),
    "IV-MVB": (0, (10, 4)),
}

METHOD_MARKERSIZE = {
    "2SPS": 6.6,
    "2SRI": 9.6,
    "GMM": 6.6,
    "IV-MVB": 6.6,
}

BIAS_Y_LIMITS = (-10, 10)
BIAS_Y_TICKS = [-10, -5, -2, -1, 0, 1, 2, 5, 10]

plt.rcParams.update({
    "font.family": "Arial",
    "font.size": 11,
    "axes.labelsize": 12,
    "axes.titlesize": 12,
    "xtick.labelsize": 10.5,
    "ytick.labelsize": 10.5,
    "figure.dpi": 150,
    "savefig.dpi": 600,
    "axes.linewidth": 1.0,
})


# ============================================================
# Utility functions
# ============================================================

def normalize_iv_strength(value):
    text = str(value).strip().lower()
    text = text.replace("-", "_")
    text = re.sub(r"\s+", "_", text)
    text = text.replace("_ivs", "").replace("_iv", "")
    text = text.replace("ivs", "").replace("iv", "")
    text = text.strip("_")

    if text in ["very_weak", "veryweak"]:
        return "very_weak"
    if text == "weak":
        return "weak"
    if text == "moderate":
        return "moderate"
    if text == "strong":
        return "strong"
    if text in ["very_strong", "verystrong"]:
        return "very_strong"

    return text


def enforce_method_order(df):
    df = df[df["method"].isin(METHOD_ORDER)].copy()
    df["method"] = pd.Categorical(
        df["method"],
        categories=METHOD_ORDER,
        ordered=True,
    )
    return df


def filter_common_setting(df, scenario, b1=None):
    df = df.copy()
    df["iv_strength_clean"] = df["iv_strength"].apply(normalize_iv_strength)

    mask = (
        (df["scenario"] == scenario) &
        (df["n"] == TARGET_N) &
        (df["iv_strength_clean"] == TARGET_IV_STRENGTH) &
        (df["c_x"].isin(CONFOUNDING_ORDER))
    )

    if b1 is not None and "b1" in df.columns:
        mask = mask & (df["b1"] == b1)

    df = df.loc[mask].copy()
    df = enforce_method_order(df)

    return df


def standard_axis_format(ax):
    ax.set_facecolor("white")

    for spine in ax.spines.values():
        spine.set_linewidth(1.0)
        spine.set_color("#222222")

    ax.tick_params(
        axis="both",
        direction="out",
        length=5,
        width=1.0,
        colors="black",
    )

    ax.grid(
        axis="y",
        color="#E6E6E6",
        linewidth=0.8,
        linestyle="-",
        zorder=0,
    )


# ============================================================
# Data loading
# ============================================================

def load_bias_data():
    df = pd.read_csv(BIAS_CSV_PATH)

    expected_cols = {
        "scenario", "b1", "n", "c_x", "iv_strength", "method",
        "min_bias", "q1_bias", "median_bias", "q3_bias", "max_bias",
    }

    missing = expected_cols - set(df.columns)
    if missing:
        raise ValueError(f"Bias CSV is missing required columns: {missing}")

    return filter_common_setting(df, SCENARIO_NULL, B1_NULL)


def load_coverage_data():
    df = pd.read_csv(COVERAGE_CSV_PATH)

    expected_cols = {
        "scenario", "b1", "n", "c_x", "iv_strength",
        "method", "n_sim", "coverage",
    }

    missing = expected_cols - set(df.columns)
    if missing:
        raise ValueError(f"Coverage CSV is missing required columns: {missing}")

    if "coverage_pct" not in df.columns:
        df["coverage_pct"] = 100 * df["coverage"]

    return filter_common_setting(df, SCENARIO_NULL, B1_NULL)


def load_type1_data():
    df = pd.read_csv(TYPE1_CSV_PATH)

    expected_cols = {
        "scenario", "n", "c_x", "iv_strength",
        "method", "n_sim", "type1_error", "type1_error_pct",
    }

    missing = expected_cols - set(df.columns)
    if missing:
        raise ValueError(f"Type I error CSV is missing required columns: {missing}")

    return filter_common_setting(df, SCENARIO_NULL, None)


def load_power_data():
    df = pd.read_csv(POWER_CSV_PATH)

    expected_cols = {
        "scenario", "b1", "n", "c_x", "iv_strength",
        "method", "n_sim", "power", "power_pct",
    }

    missing = expected_cols - set(df.columns)
    if missing:
        raise ValueError(f"Power CSV is missing required columns: {missing}")

    return filter_common_setting(df, SCENARIO_POWER, B1_POWER)


# ============================================================
# Bias boxplot helpers
# ============================================================

def visible_interval(y0, y1):
    low = max(min(y0, y1), BIAS_Y_LIMITS[0])
    high = min(max(y0, y1), BIAS_Y_LIMITS[1])

    if high <= low:
        return None

    return low, high


def draw_box_from_summary(
    ax,
    x_center,
    min_v,
    q1,
    med,
    q3,
    max_v,
    method,
    box_width,
):
    fill = METHOD_FILL[method]
    edge = METHOD_EDGE[method]

    box_interval = visible_interval(q1, q3)
    if box_interval is not None:
        box_low, box_high = box_interval
        rect = Rectangle(
            (x_center - box_width / 2, box_low),
            box_width,
            box_high - box_low,
            facecolor=fill,
            edgecolor=edge,
            linewidth=1.0,
            zorder=3,
        )
        ax.add_patch(rect)

    if BIAS_Y_LIMITS[0] <= med <= BIAS_Y_LIMITS[1]:
        ax.plot(
            [x_center - box_width / 2, x_center + box_width / 2],
            [med, med],
            color="#E41A1C",
            linewidth=2.4,
            zorder=5,
        )

    lower_whisker = visible_interval(min_v, q1)
    if lower_whisker is not None:
        low, high = lower_whisker
        ax.plot([x_center, x_center], [low, high], color=edge, linewidth=1.0, zorder=2)
        ax.plot(
            [x_center - box_width * 0.30, x_center + box_width * 0.30],
            [low, low],
            color=edge,
            linewidth=1.0,
            zorder=2,
        )

    upper_whisker = visible_interval(q3, max_v)
    if upper_whisker is not None:
        low, high = upper_whisker
        ax.plot([x_center, x_center], [low, high], color=edge, linewidth=1.0, zorder=2)
        ax.plot(
            [x_center - box_width * 0.30, x_center + box_width * 0.30],
            [high, high],
            color=edge,
            linewidth=1.0,
            zorder=2,
        )


# ============================================================
# Plotting functions
# ============================================================

def plot_bias_panel(ax, df):
    x_positions = np.arange(len(CONFOUNDING_ORDER))
    group_width = 0.65
    box_width = 0.13
    offsets = np.linspace(-group_width / 2, group_width / 2, len(METHOD_ORDER))

    for i, conf in enumerate(CONFOUNDING_ORDER):
        for method, offset in zip(METHOD_ORDER, offsets):
            row = df[
                np.isclose(df["c_x"], conf) &
                (df["method"] == method)
            ]

            if row.empty:
                continue

            row = row.iloc[0]

            draw_box_from_summary(
                ax=ax,
                x_center=i + offset,
                min_v=row["min_bias"],
                q1=row["q1_bias"],
                med=row["median_bias"],
                q3=row["q3_bias"],
                max_v=row["max_bias"],
                method=method,
                box_width=box_width,
            )

    ax.set_yscale("symlog", linthresh=1.0, linscale=1.0, base=10)

    ax.axhline(0, color="#2CA02C", linewidth=1.0, zorder=1)
    ax.axhline(1, color="#E2E2E2", linewidth=0.8, linestyle="--", zorder=1)
    ax.axhline(-1, color="#E2E2E2", linewidth=0.8, linestyle="--", zorder=1)

    ax.set_xlim(-0.7, len(CONFOUNDING_ORDER) - 0.3)
    ax.set_ylim(BIAS_Y_LIMITS)

    ax.yaxis.set_major_locator(FixedLocator(BIAS_Y_TICKS))
    ax.yaxis.set_major_formatter(FixedFormatter([str(t) for t in BIAS_Y_TICKS]))

    ax.set_xticks(x_positions)
    ax.set_xticklabels(CONFOUNDING_LABELS)

    ax.set_title("Bias", fontweight="bold", pad=8)
    ax.set_xlabel(r"Confounding Strength ($c_x=c_y$)", fontweight="bold")
    ax.set_ylabel("Estimated Bias", fontweight="bold")

    standard_axis_format(ax)


def plot_line_panel(
    ax,
    df,
    value_col,
    title,
    y_label,
    reference_y,
    y_limits=None,
    y_ticks=None,
):
    x_positions = np.arange(len(CONFOUNDING_ORDER))

    for method in METHOD_ORDER:
        y_values = []

        for conf in CONFOUNDING_ORDER:
            row = df[
                np.isclose(df["c_x"], conf) &
                (df["method"] == method)
            ]

            if row.empty:
                y_values.append(np.nan)
            else:
                y_values.append(float(row[value_col].iloc[0]))

        ax.plot(
            x_positions,
            y_values,
            color=METHOD_COLOR[method],
            linestyle=METHOD_LINESTYLE[method],
            marker=METHOD_MARKER[method],
            markersize=METHOD_MARKERSIZE[method],
            linewidth=1.25,
            markerfacecolor=METHOD_COLOR[method],
            markeredgecolor=METHOD_COLOR[method],
            markeredgewidth=0.65,
            solid_capstyle="round",
            dash_capstyle="round",
            dash_joinstyle="round",
            solid_joinstyle="round",
            label=method,
            zorder=3,
        )

    ax.axhline(
        reference_y,
        color="#33AA33",
        linewidth=1.1,
        linestyle="-",
        zorder=1,
    )

    ax.set_xlim(-0.3, len(CONFOUNDING_ORDER) - 0.7)

    if y_limits is not None:
        ax.set_ylim(y_limits)

    if y_ticks is not None:
        ax.yaxis.set_major_locator(FixedLocator(y_ticks))
        ax.yaxis.set_major_formatter(FixedFormatter([str(t) for t in y_ticks]))

    ax.set_xticks(x_positions)
    ax.set_xticklabels(CONFOUNDING_LABELS)

    ax.set_title(title, fontweight="bold", pad=8)
    ax.set_xlabel(r"Confounding Strength ($c_x=c_y$)", fontweight="bold")
    ax.set_ylabel(y_label, fontweight="bold")

    standard_axis_format(ax)


def create_figure5():
    bias_df = load_bias_data()
    coverage_df = load_coverage_data()
    type1_df = load_type1_data()
    power_df = load_power_data()

    fig, axes = plt.subplots(
        2,
        2,
        figsize=(11.2, 8.7),
        facecolor="white",
    )

    axes = axes.ravel()

    plot_bias_panel(axes[0], bias_df)

    plot_line_panel(
        ax=axes[1],
        df=coverage_df,
        value_col="coverage_pct",
        title="Coverage",
        y_label="Coverage Estimate (%)",
        reference_y=95,
        y_limits=(85, 101),
        y_ticks=[85, 90, 95, 100],
    )

    plot_line_panel(
        ax=axes[2],
        df=type1_df,
        value_col="type1_error_pct",
        title="Type I Error",
        y_label="Type I Error Estimate (%)",
        reference_y=5,
        y_limits=(0, 15),
        y_ticks=[0, 5, 10, 15],
    )

    plot_line_panel(
        ax=axes[3],
        df=power_df,
        value_col="power_pct",
        title="Power",
        y_label="Power Estimate (%)",
        reference_y=80,
        y_limits=(0, 100),
        y_ticks=[0, 20, 40, 60, 80, 100],
    )

    fig.subplots_adjust(
        left=0.10,
        right=0.96,
        top=0.92,
        bottom=0.11,
        wspace=0.24,
        hspace=0.38,
    )

    legend_handles = [
        plt.Line2D(
            [0],
            [0],
            color=METHOD_COLOR[m],
            linestyle=METHOD_LINESTYLE[m],
            marker=METHOD_MARKER[m],
            markersize=METHOD_MARKERSIZE[m],
            linewidth=1.25,
            markerfacecolor=METHOD_COLOR[m],
            markeredgecolor=METHOD_COLOR[m],
            markeredgewidth=0.65,
            label=m,
        )
        for m in METHOD_ORDER
    ]

    leg = fig.legend(
        handles=legend_handles,
        loc="lower center",
        ncol=4,
        frameon=False,
        handlelength=2.2,
        handletextpad=0.5,
        columnspacing=2.0,
        fontsize=11,
    )

    for text in leg.get_texts():
        text.set_fontweight("bold")

    png_path = OUTPUT_DIR / "Figure5_operating_characteristics.png"
    svg_path = OUTPUT_DIR / "Figure5_operating_characteristics.svg"

    fig.savefig(
        png_path,
        dpi=600,
        facecolor="white",
        edgecolor="white",
        bbox_inches="tight",
    )

    fig.savefig(
        svg_path,
        facecolor="white",
        edgecolor="white",
        bbox_inches="tight",
    )

    plt.close(fig)

    return png_path, svg_path


# ============================================================
# PowerPoint helpers
# ============================================================

def make_figure_title():
    title = (
        "Figure 5. Comparison of MR method performance across confounding strengths\n"
        "under weak instruments"
    )
    subtitle = (
        r"Fixed setting: n = 1000, weak IVs, and "
        r"c_x = c_y in {0.5, 1.0, 1.5}. "
        r"Bias, coverage, and Type I error are shown for Scenario A (b1 = 0); "
        r"power is shown for Scenario B (b1 = 1). "
        r"See Table 1 for more details on the parameter settings of scenarios."
    )

    return title, subtitle


def add_slide_decorations(slide, title, subtitle):
    title_box = slide.shapes.add_textbox(
        Inches(0.45),
        Inches(0.14),
        Inches(14.0),
        Inches(0.90),
    )

    tf = title_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 20, 20)

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.45),
        Inches(0.90),
        Inches(14.0),
        Inches(0.88),
    )

    tf2 = subtitle_box.text_frame
    tf2.word_wrap = True

    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.name = "Arial"
    p2.font.size = Pt(13.5)
    p2.font.color.rgb = RGBColor(60, 60, 60)


def create_powerpoint(png_path):
    prs = Presentation()
    prs.slide_width = Inches(15.0)
    prs.slide_height = Inches(11.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    title, subtitle = make_figure_title()
    add_slide_decorations(slide, title, subtitle)

    img_left = Inches(1.35)
    img_top = Inches(1.92)
    img_width = Inches(10.75)
    img_height = Inches(8.25)

    border_pad = Inches(0.04)

    border_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        img_left - border_pad,
        img_top - border_pad,
        img_width + 2 * border_pad,
        img_height + 2 * border_pad,
    )

    border_box.fill.background()
    border_box.line.color.rgb = RGBColor(55, 55, 55)
    border_box.line.width = Pt(1.5)

    slide.shapes.add_picture(
        str(png_path),
        img_left,
        img_top,
        width=img_width,
        height=img_height,
    )

    prs.save(PPTX_OUT)


# ============================================================
# Main
# ============================================================

def main():
    png_path, svg_path = create_figure5()
    create_powerpoint(png_path)

    print("\nSuccess! Figure 5 PowerPoint created successfully.\n")
    print(f"Output PPTX: {PPTX_OUT}")
    print(f"Figure 5 PNG: {png_path}")
    print(f"Figure 5 SVG: {svg_path}")


if __name__ == "__main__":
    main()