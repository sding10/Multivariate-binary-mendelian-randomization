# Save as: create_figure_4.py
# Run from the folder that contains the Data folder.

from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.ticker import FixedLocator, FixedFormatter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# File paths
# ============================================================

TYPE1_CSV_PATH = Path("Data") / "Figure_TypeIError_Summary.csv"
POWER_CSV_PATH = Path("Data") / "Figure_Power_Summary.csv"

OUTPUT_DIR = Path("Figure4_type1_power_from_csv_output")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

PPTX_OUT = OUTPUT_DIR / "Figure4_type1_power_from_csv.pptx"


# ============================================================
# Figure configuration
# ============================================================

FIGURE_CONFIG = {
    "4a": {"conf": 0.5},
    "4b": {"conf": 1.0},
    "4c": {"conf": 1.5},
}

POWER_B1_TO_KEEP = 1


# ============================================================
# Plot settings
# ============================================================

METHOD_ORDER = ["2SPS", "2SRI", "GMM", "IV-MVB"]

METHOD_COLOR = {
    "2SPS": "#CFCFCF",
    "2SRI": "#AFAFAF",
    "GMM": "#8E8E8E",
    "IV-MVB": "#111111",
}

METHOD_MARKER = {
    "2SPS": "o",
    "2SRI": "*",
    "GMM": "D",
    "IV-MVB": "s",
}

METHOD_LINESTYLE = {
    "2SPS": (0, (10, 4)),
    "2SRI": (0, (10, 4)),
    "GMM": (0, (10, 4)),
    "IV-MVB": (0, (10, 4)),
}

METHOD_MARKERSIZE = {
    "2SPS": 6.6,
    "2SRI": 9.6,
    "GMM": 6.6,
    "IV-MVB": 6.6,
}

IV_STRENGTH_ORDER = [
    "very_weak",
    "weak",
    "moderate",
    "strong",
    "very_strong",
]

IV_STRENGTH_LABEL = {
    "very_weak": "Very Weak\nIVs",
    "weak": "Weak\nIVs",
    "moderate": "Moderate\nIVs",
    "strong": "Strong\nIVs",
    "very_strong": "Very Strong\nIVs",
}

SAMPLE_SIZE_ORDER = [500, 1500, 2500, 5000, 7500, 10000]

TYPE1_Y_LIMITS = (0, 15)
TYPE1_Y_TICKS = [0, 5, 10, 15]

POWER_Y_LIMITS = (0, 100)
POWER_Y_TICKS = [0, 20, 40, 60, 80, 100]

plt.rcParams.update({
    "font.family": "Arial",
    "font.size": 11,
    "axes.labelsize": 12,
    "axes.titlesize": 12,
    "xtick.labelsize": 10.5,
    "ytick.labelsize": 10.5,
    "figure.dpi": 150,
    "savefig.dpi": 600,
    "axes.linewidth": 1.0,
})


# ============================================================
# Data loading
# ============================================================

def load_type1_csv(csv_path):
    df = pd.read_csv(csv_path)

    expected_cols = {
        "scenario",
        "n",
        "c_x",
        "iv_strength",
        "method",
        "n_sim",
        "type1_error",
        "type1_error_pct",
    }

    missing = expected_cols - set(df.columns)
    if missing:
        raise ValueError(f"Type I error CSV is missing required columns: {missing}")

    df = df[df["scenario"].isin(["A", "C"])].copy()
    df = df[df["method"].isin(METHOD_ORDER)].copy()

    df["metric"] = "Type I Error"
    df["metric_pct"] = df["type1_error_pct"]

    return df


def load_power_csv(csv_path):
    df = pd.read_csv(csv_path)

    expected_cols = {
        "scenario",
        "b1",
        "n",
        "c_x",
        "iv_strength",
        "method",
        "n_sim",
        "power",
        "power_pct",
    }

    missing = expected_cols - set(df.columns)
    if missing:
        raise ValueError(f"Power CSV is missing required columns: {missing}")

    df = df[df["scenario"].isin(["B", "D"])].copy()

    # Match the Figure 3 setup: keep beta1 = 1 and ignore beta1 = 2 or 3.
    df = df[df["b1"] == POWER_B1_TO_KEEP].copy()

    df = df[df["method"].isin(METHOD_ORDER)].copy()

    df["metric"] = "Power"
    df["metric_pct"] = df["power_pct"]

    return df


def load_and_combine_data():
    df_type1 = load_type1_csv(TYPE1_CSV_PATH)
    df_power = load_power_csv(POWER_CSV_PATH)

    df_all = pd.concat([df_type1, df_power], ignore_index=True)

    df_all["method"] = pd.Categorical(
        df_all["method"],
        categories=METHOD_ORDER,
        ordered=True,
    )

    df_all["iv_strength"] = df_all["iv_strength"].astype(str).str.strip()

    def choose_x(row):
        if row["scenario"] in ["A", "B"]:
            return row["iv_strength"]
        return int(row["n"])

    df_all["x_value"] = df_all.apply(choose_x, axis=1)

    return df_all


def get_figure_data(df_all, conf_strength):
    return df_all[np.isclose(df_all["c_x"], conf_strength)].copy()


# ============================================================
# Plotting
# ============================================================

def plot_panel(ax, df, scenario, panel_label):
    panel = df[df["scenario"] == scenario].copy()

    if scenario in ["A", "B"]:
        x_order = IV_STRENGTH_ORDER
        x_label = "Instrument Strength"
    else:
        x_order = SAMPLE_SIZE_ORDER
        x_label = "Sample Size"

    if scenario in ["A", "C"]:
        y_label = "Type I Error Estimate (%)"
        y_limits = TYPE1_Y_LIMITS
        y_ticks = TYPE1_Y_TICKS
        reference_y = 5
    else:
        y_label = "Power Estimate (%)"
        y_limits = POWER_Y_LIMITS
        y_ticks = POWER_Y_TICKS
        reference_y = 80

    x_positions = np.arange(len(x_order))

    for method in METHOD_ORDER:
        y_values = []

        for x_val in x_order:
            row = panel[
                (panel["x_value"] == x_val) &
                (panel["method"] == method)
            ]

            if row.empty:
                y_values.append(np.nan)
            else:
                y_values.append(float(row["metric_pct"].iloc[0]))

        ax.plot(
            x_positions,
            y_values,
            color=METHOD_COLOR[method],
            linestyle=METHOD_LINESTYLE[method],
            marker=METHOD_MARKER[method],
            markersize=METHOD_MARKERSIZE[method],
            linewidth=1.25,
            markerfacecolor=METHOD_COLOR[method],
            markeredgecolor=METHOD_COLOR[method],
            markeredgewidth=0.65,
            solid_capstyle="round",
            dash_capstyle="round",
            dash_joinstyle="round",
            solid_joinstyle="round",
            label=method,
            zorder=3,
        )

    ax.axhline(
        reference_y,
        color="#33AA33",
        linewidth=1.1,
        linestyle="-",
        zorder=1,
    )

    ax.set_xlim(-0.3, len(x_order) - 0.7)
    ax.set_ylim(y_limits)

    ax.yaxis.set_major_locator(FixedLocator(y_ticks))
    ax.yaxis.set_major_formatter(FixedFormatter([str(t) for t in y_ticks]))

    ax.set_xticks(x_positions)

    if scenario in ["A", "B"]:
        ax.set_xticklabels([IV_STRENGTH_LABEL[x] for x in x_order])
    else:
        ax.set_xticklabels([str(x) for x in x_order])

    ax.set_ylabel(y_label, fontweight="bold")
    ax.set_xlabel(x_label, fontweight="bold")

    ax.set_facecolor("white")

    for spine in ax.spines.values():
        spine.set_linewidth(1.0)
        spine.set_color("#222222")

    ax.tick_params(
        axis="both",
        direction="out",
        length=5,
        width=1.0,
        colors="black",
    )

    ax.grid(axis="y", color="#E6E6E6", linewidth=0.8, linestyle="-", zorder=0)

    ax.text(
        0.02,
        0.97,
        panel_label,
        transform=ax.transAxes,
        ha="left",
        va="top",
        fontsize=20,
        fontweight="bold",
        color="black",
    )


def create_one_figure(df, fig_label, conf_strength):
    fig, axes = plt.subplots(
        2,
        2,
        figsize=(11.2, 8.7),
        facecolor="white",
    )

    axes = axes.ravel()

    plot_panel(axes[0], df, "A", "A")
    plot_panel(axes[1], df, "B", "B")
    plot_panel(axes[2], df, "C", "C")
    plot_panel(axes[3], df, "D", "D")

    fig.subplots_adjust(
        left=0.10,
        right=0.96,
        top=0.96,
        bottom=0.11,
        wspace=0.24,
        hspace=0.32,
    )

    legend_handles = [
        plt.Line2D(
            [0],
            [0],
            color=METHOD_COLOR[m],
            linestyle=METHOD_LINESTYLE[m],
            marker=METHOD_MARKER[m],
            markersize=METHOD_MARKERSIZE[m],
            linewidth=1.25,
            markerfacecolor=METHOD_COLOR[m],
            markeredgecolor=METHOD_COLOR[m],
            markeredgewidth=0.65,
            solid_capstyle="round",
            dash_capstyle="round",
            dash_joinstyle="round",
            solid_joinstyle="round",
            label=m,
        )
        for m in METHOD_ORDER
    ]

    leg = fig.legend(
        handles=legend_handles,
        loc="lower center",
        ncol=4,
        frameon=False,
        handlelength=2.2,
        handletextpad=0.5,
        columnspacing=2.0,
        fontsize=11,
    )

    for text in leg.get_texts():
        text.set_fontweight("bold")

    conf_tag = str(conf_strength).replace(".", "p")

    png_path = OUTPUT_DIR / f"Figure_{fig_label}_type1_power_c{conf_tag}.png"
    svg_path = OUTPUT_DIR / f"Figure_{fig_label}_type1_power_c{conf_tag}.svg"

    fig.savefig(
        png_path,
        dpi=600,
        facecolor="white",
        edgecolor="white",
        bbox_inches="tight",
    )

    fig.savefig(
        svg_path,
        facecolor="white",
        edgecolor="white",
        bbox_inches="tight",
    )

    plt.close(fig)

    return png_path, svg_path


# ============================================================
# PowerPoint helpers
# ============================================================

def make_figure_title(fig_label, conf):
    title = (
        f"Figure {fig_label}. Type I error and power estimates for each MR method\n"
        f"in each simulation scenario with 30 instruments"
    )

    subtitle = (
        f"Confounding strength: c_x = c_y = {conf:.1f}. "
        f"Scenarios A and C show Type I error under beta1 = 0; "
        f"Scenarios B and D show power under beta1 = {POWER_B1_TO_KEEP}. "
        f"Green reference lines indicate nominal 5% Type I error in A/C "
        f"and an 80% power benchmark in B/D. "
        f"See Table 1 for more details on the parameter settings in each scenario."
    )

    return title, subtitle


def add_slide_decorations(slide, title, subtitle):
    title_box = slide.shapes.add_textbox(
        Inches(0.45),
        Inches(0.16),
        Inches(14.0),
        Inches(0.85),
    )

    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 20, 20)

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.45),
        Inches(0.88),
        Inches(14.0),
        Inches(0.80),
    )

    tf2 = subtitle_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.name = "Arial"
    p2.font.size = Pt(13.5)
    p2.font.color.rgb = RGBColor(60, 60, 60)


def create_powerpoint(image_records):
    prs = Presentation()
    prs.slide_width = Inches(15.0)
    prs.slide_height = Inches(11.5)

    blank_layout = prs.slide_layouts[6]

    for rec in image_records:
        slide = prs.slides.add_slide(blank_layout)

        title, subtitle = make_figure_title(rec["fig_label"], rec["conf"])
        add_slide_decorations(slide, title, subtitle)

        img_left = Inches(1.25)
        img_top = Inches(1.90)
        img_width = Inches(10.95)
        img_height = Inches(8.25)

        border_pad = Inches(0.04)

        border_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            img_left - border_pad,
            img_top - border_pad,
            img_width + 2 * border_pad,
            img_height + 2 * border_pad,
        )

        border_box.fill.background()
        border_box.line.color.rgb = RGBColor(55, 55, 55)
        border_box.line.width = Pt(1.5)

        slide.shapes.add_picture(
            str(rec["png"]),
            img_left,
            img_top,
            width=img_width,
            height=img_height,
        )

    prs.save(PPTX_OUT)


# ============================================================
# Main
# ============================================================

def main():
    df_all = load_and_combine_data()

    image_records = []

    for fig_label, cfg in FIGURE_CONFIG.items():
        conf = cfg["conf"]
        df_fig = get_figure_data(df_all, conf)

        png_path, svg_path = create_one_figure(
            df=df_fig,
            fig_label=fig_label,
            conf_strength=conf,
        )

        image_records.append({
            "fig_label": fig_label,
            "conf": conf,
            "png": png_path,
            "svg": svg_path,
        })

    create_powerpoint(image_records)

    print("\nSuccess! Figure 4 PowerPoint created successfully.\n")
    print(f"Output PPTX: {PPTX_OUT}\n")

    for rec in image_records:
        print(f"Figure {rec['fig_label']} PNG: {rec['png']}")
        print(f"Figure {rec['fig_label']} SVG: {rec['svg']}")


if __name__ == "__main__":
    main()