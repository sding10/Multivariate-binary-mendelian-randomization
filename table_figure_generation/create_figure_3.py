# ============================================================
# Create Figure 3a, 3b, 3c from raw CSV summary file
# Coverage estimates for each MR method in each simulation scenario
# with 30 instruments
# ============================================================

from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.patches import Patch
from matplotlib.ticker import FixedLocator, FixedFormatter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# File paths
# ============================================================

CSV_PATH = Path("Data") / "Figure_Coverage_Summary.csv"

OUTPUT_DIR = Path("Figure3_coverage_from_csv_output")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

PPTX_OUT = OUTPUT_DIR / "Figure3_coverage_from_csv.pptx"


# ============================================================
# Figure configuration
# ============================================================

FIGURE_CONFIG = {
    "3a": {"conf": 0.5},
    "3b": {"conf": 1.0},
    "3c": {"conf": 1.5},
}


# ============================================================
# Plot settings
# ============================================================

METHOD_ORDER = ["2SPS", "2SRI", "GMM", "IV-MVB"]

METHOD_COLOR = {
    "2SPS": "#CFCFCF",   # very light gray
    "2SRI": "#AFAFAF",   # medium-light gray
    "GMM":  "#8E8E8E",   # darker gray
    "IV-MVB": "#111111", # near-black
}

METHOD_MARKER = {
    "2SPS": "o",   # circle
    "2SRI": "*",   # star
    "GMM": "D",    # diamond
    "IV-MVB": "s", # square
}

# Long elegant dash pattern similar to the reference image
METHOD_LINESTYLE = {
    "2SPS": (0, (10, 4)),
    "2SRI": (0, (10, 4)),
    "GMM":  (0, (10, 4)),
    "IV-MVB": (0, (10, 4)),
}

# Optional: separate marker sizes for a more polished look
METHOD_MARKERSIZE = {
    "2SPS": 6.6,
    "2SRI": 9.6,
    "GMM":  6.6,
    "IV-MVB": 6.6,
}

IV_STRENGTH_ORDER = [
    "Very Weak IVs",
    "Weak IVs",
    "Moderate IVs",
    "Strong IVs",
    "Very Strong IVs",
]

SAMPLE_SIZE_ORDER = [500, 1500, 2500, 5000, 7500, 10000]

Y_LIMITS = (85, 101)
Y_TICKS = [85, 90, 95, 100]

plt.rcParams.update({
    "font.family": "Arial",
    "font.size": 11,
    "axes.labelsize": 12,
    "axes.titlesize": 12,
    "xtick.labelsize": 10.5,
    "ytick.labelsize": 10.5,
    "figure.dpi": 150,
    "savefig.dpi": 600,
    "axes.linewidth": 1.0,
})


# ============================================================
# Data loading and filtering
# ============================================================

def clean_x_label(value):
    text = str(value)

    replacements = {
        "Very Weak IVs": "Very Weak\nIVs",
        "Weak IVs": "Weak\nIVs",
        "Moderate IVs": "Moderate\nIVs",
        "Strong IVs": "Strong\nIVs",
        "Very Strong IVs": "Very Strong\nIVs",
    }

    return replacements.get(text, text)


def load_and_filter_raw_csv(csv_path):
    df = pd.read_csv(csv_path)

    expected_cols = {
        "scenario",
        "b1",
        "n",
        "c_x",
        "iv_strength",
        "method",
        "n_sim",
        "coverage",
    }

    missing = expected_cols - set(df.columns)
    if missing:
        raise ValueError(f"CSV is missing required columns: {missing}")

    # Create coverage percentage if needed.
    if "coverage_pct" not in df.columns:
        df["coverage_pct"] = 100 * df["coverage"]

    # Required filtering:
    # Ignore b1 = 2 and b1 = 3.
    # Keep only the intended scenario-specific b1 values.
    mask = (
        ((df["scenario"] == "A") & (df["b1"] == 0)) |
        ((df["scenario"] == "B") & (df["b1"] == 1)) |
        ((df["scenario"] == "C") & (df["b1"] == 0)) |
        ((df["scenario"] == "D") & (df["b1"] == 1))
    )

    df = df.loc[mask].copy()

    df = df[df["method"].isin(METHOD_ORDER)].copy()
    df["method"] = pd.Categorical(
        df["method"],
        categories=METHOD_ORDER,
        ordered=True
    )

    def choose_x(row):
        if row["scenario"] in ["A", "B"]:
            return row["iv_strength"]
        return int(row["n"])

    df["x_value"] = df.apply(choose_x, axis=1)

    return df


def get_figure_data(df_all, conf_strength):
    return df_all[np.isclose(df_all["c_x"], conf_strength)].copy()


# ============================================================
# Plotting
# ============================================================

def plot_panel(ax, df, scenario, panel_label):
    panel = df[df["scenario"] == scenario].copy()

    if scenario in ["A", "B"]:
        x_order = IV_STRENGTH_ORDER
        x_label = "Instrument Strength"
    else:
        x_order = SAMPLE_SIZE_ORDER
        x_label = "Sample Size"

    x_positions = np.arange(len(x_order))

    for method in METHOD_ORDER:
        y_values = []

        for x_val in x_order:
            row = panel[
                (panel["x_value"] == x_val) &
                (panel["method"] == method)
            ]

            if row.empty:
                y_values.append(np.nan)
            else:
                y_values.append(float(row["coverage_pct"].iloc[0]))

        ax.plot(
            x_positions,
            y_values,
            color=METHOD_COLOR[method],
            linestyle=METHOD_LINESTYLE[method],
            marker=METHOD_MARKER[method],
            markersize=METHOD_MARKERSIZE[method],
            linewidth=1.25,
            markerfacecolor=METHOD_COLOR[method],   # filled markers like reference
            markeredgecolor=METHOD_COLOR[method],
            markeredgewidth=0.65,
            solid_capstyle="round",
            dash_capstyle="round",
            dash_joinstyle="round",
            solid_joinstyle="round",
            label=method,
            zorder=3,
        )

    ax.axhline(
        95,
        color="#33AA33",
        linewidth=1.1,
        linestyle="-",
        zorder=1,
    )


    ax.set_xlim(-0.3, len(x_order) - 0.7)
    ax.set_ylim(Y_LIMITS)

    ax.yaxis.set_major_locator(FixedLocator(Y_TICKS))
    ax.yaxis.set_major_formatter(FixedFormatter([str(t) for t in Y_TICKS]))

    ax.set_xticks(x_positions)
    ax.set_xticklabels([clean_x_label(x) for x in x_order])

    ax.set_ylabel("Coverage Estimate (%)", fontweight="bold")
    ax.set_xlabel(x_label, fontweight="bold")

    ax.set_facecolor("white")

    for spine in ax.spines.values():
        spine.set_linewidth(1.0)
        spine.set_color("#222222")

    ax.tick_params(
        axis="both",
        direction="out",
        length=5,
        width=1.0,
        colors="black"
    )

    ax.grid(axis="y", color="#E6E6E6", linewidth=0.8, linestyle="-", zorder=0)

    ax.text(
        0.02,
        0.97,
        panel_label,
        transform=ax.transAxes,
        ha="left",
        va="top",
        fontsize=20,
        fontweight="bold",
        color="black",
    )


def create_one_figure(df, fig_label, conf_strength):
    fig, axes = plt.subplots(
        2,
        2,
        figsize=(11.2, 8.7),
        facecolor="white"
    )

    axes = axes.ravel()

    plot_panel(axes[0], df, "A", "A")
    plot_panel(axes[1], df, "B", "B")
    plot_panel(axes[2], df, "C", "C")
    plot_panel(axes[3], df, "D", "D")

    fig.subplots_adjust(
        left=0.10,
        right=0.96,
        top=0.96,
        bottom=0.11,
        wspace=0.24,
        hspace=0.32,
    )
    legend_handles = [
        plt.Line2D(
            [0],
            [0],
            color=METHOD_COLOR[m],
            linestyle=METHOD_LINESTYLE[m],
            marker=METHOD_MARKER[m],
            markersize=METHOD_MARKERSIZE[m],
            linewidth=1.25,
            markerfacecolor=METHOD_COLOR[m],
            markeredgecolor=METHOD_COLOR[m],
            markeredgewidth=0.65,
            solid_capstyle="round",
            dash_capstyle="round",
            dash_joinstyle="round",
            solid_joinstyle="round",
            label=m,
        )
        for m in METHOD_ORDER
    ]

    leg = fig.legend(
        handles=legend_handles,
        loc="lower center",
        ncol=4,
        frameon=False,
        handlelength=2.2,
        handletextpad=0.5,
        columnspacing=2.0,
        fontsize=11,
    )

    for text in leg.get_texts():
        text.set_fontweight("bold")

    conf_tag = str(conf_strength).replace(".", "p")

    png_path = OUTPUT_DIR / f"Figure_{fig_label}_coverage_c{conf_tag}.png"
    svg_path = OUTPUT_DIR / f"Figure_{fig_label}_coverage_c{conf_tag}.svg"

    fig.savefig(
        png_path,
        dpi=600,
        facecolor="white",
        edgecolor="white",
        bbox_inches="tight",
    )

    fig.savefig(
        svg_path,
        facecolor="white",
        edgecolor="white",
        bbox_inches="tight",
    )

    plt.close(fig)

    return png_path, svg_path


# ============================================================
# PowerPoint helpers
# ============================================================
def make_figure_title(fig_label, conf):
    title = (
        f"Figure {fig_label}. Coverage estimates for each MR method\n"
        f"in each simulation scenario with 30 instruments"
    )

    subtitle = (
        f"Confounding strength: c_x = c_y = {conf:.1f}. "
        f"The dashed horizontal line indicates nominal 95% coverage.\n"
        f"See Table 1 for more details on the parameter settings in each scenario."
    )

    return title, subtitle
    
def make_figure_title(fig_label, conf):
    title = (
        f"Figure {fig_label}. Coverage estimates for each MR method "
        f"in each simulation scenario with 30 instruments"
    )
    subtitle = (
        f"Confounding strength: c_x = c_y = {conf:.1f}. "
        f"The dashed horizontal line indicates nominal 95% coverage. "
        f"See Table 1 for more details on the parameter settings in each scenario."
    )

    return title, subtitle


def add_slide_decorations(slide, title, subtitle):
    title_box = slide.shapes.add_textbox(
        Inches(0.45),
        Inches(0.18),
        Inches(14.0),
        Inches(0.70)
    )

    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 20, 20)

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.45),
        Inches(0.65),
        Inches(14.0),
        Inches(0.70)
    )

    tf2 = subtitle_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.name = "Arial"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(60, 60, 60)


def create_powerpoint(image_records):
    prs = Presentation()
    prs.slide_width = Inches(15.0)
    prs.slide_height = Inches(11.5)

    blank_layout = prs.slide_layouts[6]

    for rec in image_records:
        slide = prs.slides.add_slide(blank_layout)

        title, subtitle = make_figure_title(rec["fig_label"], rec["conf"])
        add_slide_decorations(slide, title, subtitle)

        img_left = Inches(1.35)
        img_top = Inches(1.78)
        img_width = Inches(10.75)
        img_height = Inches(8.35)

        border_pad = Inches(0.04)

        border_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            img_left - border_pad,
            img_top - border_pad,
            img_width + 2 * border_pad,
            img_height + 2 * border_pad,
        )

        border_box.fill.background()
        border_box.line.color.rgb = RGBColor(55, 55, 55)
        border_box.line.width = Pt(1.5)

        slide.shapes.add_picture(
            str(rec["png"]),
            img_left,
            img_top,
            width=img_width,
            height=img_height,
        )

    prs.save(PPTX_OUT)


# ============================================================
# Main
# ============================================================

def main():
    df_all = load_and_filter_raw_csv(CSV_PATH)

    image_records = []

    for fig_label, cfg in FIGURE_CONFIG.items():
        conf = cfg["conf"]
        df_fig = get_figure_data(df_all, conf)

        png_path, svg_path = create_one_figure(
            df=df_fig,
            fig_label=fig_label,
            conf_strength=conf,
        )

        image_records.append({
            "fig_label": fig_label,
            "conf": conf,
            "png": png_path,
            "svg": svg_path,
        })

    create_powerpoint(image_records)

    print("\nSuccess! Figure 3 PowerPoint created successfully.\n")
    print(f"Output PPTX: {PPTX_OUT}\n")

    for rec in image_records:
        print(f"Figure {rec['fig_label']} PNG: {rec['png']}")
        print(f"Figure {rec['fig_label']} SVG: {rec['svg']}")


if __name__ == "__main__":
    main()